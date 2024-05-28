from PIL._imaging import display
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import requests
import codecs
import csv
import time
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import numpy as np
from Framaform import FramaFormResult #formulaire avec nom prénom mention dominante photo legende, csv
from urllib.request import urlopen, urlretrieve
from PIL import Image
from PIL import ImageOps
from PIL import ImageDraw
import os
import glob
import urllib.request
import urllib
from io import BytesIO


# Start a session so we can have persistant cookies
session = requests.session()
# ctrl shift i
# copy/paste cookie from the web page connexion
cookie_frama = {
    'SSESSc4ca4917a9afea424a5a393168053720':'uKFDSAd8EBRHhq73nPi1ioNdEUTybDuCqkFDgM1nu2A'
} #navigateur connexion cookie de la session

# place the csv file into the folder
csv_path = 'rdd.csv'

# dict_mention is a dictonnary with an id for each mention
dict_mention = {'71': 'GSI - CTE', '73': 'GSI - SCOM', '72': 'GSI - DS', '82': 'CVT - AET', '81': 'CVT - SIC',
                '31': 'INFO - IA', '32': 'INFO - SL', '33': 'INFO - ASI', '34': 'INFO - CYBER', '35': 'INFO - MSc AI',
                '53': 'ENE - E2', '52': 'ENE - PEG', '51': 'ENE - RE', '54': 'ENE - SES', '61': 'VSE - ESP',
                '62': 'VSE - HSB', '41': 'SCOC - SRI', '43': 'SCOC - MACS', '42': 'SCOC - OCENE', '22': 'PNT - PSY',
                '21': 'PNT - QTE', '11': 'MDS - MMF', '13': 'MDS - SDI-PS', '12': 'MDS - SDI-M', '14': 'MDS - MSc DSBA'}

# dict_dominante is a dictonnary with an id for each dominante
dict_dominante = {'7': 'GSI', '8': 'CVT', '3': 'INFO', '5': 'ENE', '6': 'VSE', '4': 'SCOC', '2': 'PNT', '1': 'MDS'}


def add_citation(shapes, text_cit):
    """
    to add citation to the slide
    :param shapes: shape
    :param text_cit:
    :return:
    """
    left = Inches(0.1)
    top = Inches(5.8)
    width = Inches(16)
    height = Inches(1)
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text_cit
    font = run.font
    font.name = 'OPEN SANS'
    font.size = Pt(20)
    # p.alignment = PP_ALIGN.JUSTIFY
    font.color.rgb = RGBColor(255, 255, 255)
    fill = shape.fill
    fill.background()
    line = shape.line
    line.fill.background()


def add_title(shapes, surname, name):
    left = Inches(0.1)
    top = Inches(5.1)
    width = Inches(16)
    height = Inches(1)
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = surname + ' ' + name.upper()
    font = run.font
    font.name = 'OPEN SANS'
    font.size = Pt(36)
    # p.alignment = PP_ALIGN.JUSTIFY
    font.color.rgb = RGBColor(255, 255, 255)
    fill = shape.fill
    fill.background()
    line = shape.line
    line.fill.background()

def load_img(url_image, cookie_frama, session):
    # Download image
    file_path = r'test.png'
    with open(file_path, 'wb') as handle:
        response = session.get(url_image, cookies=cookie_frama)
        if not response.ok:
            print("Something went wrong")
        else:
            for block in response.iter_content(1024):
                handle.write(block)
    image = Image.open(file_path)
    image = ImageOps.exif_transpose(image)
    image.save(file_path)
    return file_path


def wait_for_file(filepath):
    wait_time = 1
    while is_locked(filepath):
        time.sleep(wait_time)
    #handle.close()


def add_image(shapes, path_img, SLIDE_WIDTH, SLIDE_HEIGHT):
    toph = Inches(1)
    #widthimg , heightimg = img.size
    #aspect_ratio = heightimg / widthimg
    # Hauteur voulue en pixel
    desired_size = 1000
    img = Image.open(path_img)
    img=img.convert('RGBA')
    # On veut faire une image carrée de hauteur 1000px -> on
    # prend la plus petit dimension de l'image et on la ramène à 1000px
    # 1er cas : la largeur est plus grande que la hauteur : on ramène la hauteur à
    # la hauteur désirée et on ajuste la largeur
    if img.size[0] >= img.size[1]:
        ratio_old_to_new = desired_size / img.size[1]
        img = img.resize((int(ratio_old_to_new * img.size[0]), desired_size))
        # Coordonnées de l'ellipse
        # Point en haut à gauche
        x_0, y_0 = int(img.size[0] / 2 - desired_size / 2), 0
        # Point en bas à droite de l'ellipse
        x_1, y_1 = int(desired_size / 2 + img.size[0] / 2), desired_size

    # Sinon on ramène la largeur à 1000 et on ajuste la hauteur
    else:
        ratio_old_to_new = desired_size / img.size[0]
        img = img.resize((desired_size, int(ratio_old_to_new * img.size[1])))
        # Coordonnées de l'ellipse
        # Point en haut à gauche
        x_0, y_0 = 0, int(img.size[1] / 2 - desired_size / 2)
        # Point en bas à droite de l'ellipse
        x_1, y_1 = desired_size, int(desired_size / 2 + img.size[1] / 2)

    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)

    # # Crop parant d'en haute à gauche
    # max_size = min(img.size[0], img.size[1])
    # draw.ellipse((0, 0, max_size, max_size), fill=255)
    # # Crop centré
    draw.ellipse((x_0, y_0, x_1, y_1), fill=255)

    # Apply the mask to the image
    img.putalpha(mask)

    # Suppress the image outside the mask
    img = img.crop(img.getbbox())
    new_width=img.size[0]
    new_width=new_width*3700
    new_height=img.size[1]
    new_height=new_height*3700
    img.save("photo.png")
    #left= (SLIDE_WIDTH - img.size[0])/2
    left=Inches(5.9)
    shapes.add_picture("photo.png", left,toph,new_width,new_height)



    #if widthimg <= heightimg:
        # Calculate the size of the displayed image
     #       img_display_height = height
      #      img_display_width = img_display_height / aspect_ratio
       #     img=img.resize((int(height/aspect_ratio),height))


        # lefth = (Inches(4) + SLIDE_WIDTH - img_display_width) / 2
            #lefth = Inches(7) + (Inches(16) - Inches(7) - img_display_width) / 2
            #shape_pic = shapes.add_picture(path_img, lefth, toph, width = img_diplaheight=height)
            #shape_pic = shapes.add_picture(path_img, Inches(6.3), Inches(1.57),
                              # width= Inches(3.93), height= Inches(3.15))


   # else:
        # Calculate the size of the displayed image
    #        img_display_width = width
     #       img_display_height=img_display_width*aspect_ratio
      #      img=img.resize((width,int(width*aspect_ratio)))
        #img_display_height = aspect_ratio * img_display_width
            #topw=img_display_height
        #topw = (SLIDE_HEIGHT - img_display_height) / 2
            #shape_pic = shapes.add_picture(path_img, left, topw, width=width)
            #shape_pic = shapes.add_picture(path_img, Inches(5.90), Inches(1.57),
              #                 width= Inches(4.72), height= Inches(3.15))
    #left= (SLIDE_WIDTH - img_display_width)/2
    #mask = Image.new("L", img.size, 0)
    #draw = ImageDraw.Draw(mask)
    #max_size = min(img.size[0], img.size[1])
    #draw.ellipse((0, 0, max_size, max_size), fill=255)

    # Apply the mask to the image
    #img.putalpha(mask)

    # Suppress the image outside the mask
    #img = img.crop(img.getbbox())
    #with BytesIO() as output:
     #   img.save(output, format="PNG")
        #shape_pic = shapes.add_picture(img_path, left, toph, width = img_display_width,height=img_display_height)


def add_slide_dominante(dominante_name, prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    add_title(shapes, dominante_name, '')



def add_mention(shapes, mention_eleve):
    left = Inches(0.1)
    top = Inches(8)
    width = Inches(16)
    height = Inches(1)
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = ' ' + mention_eleve.upper()
    font = run.font
    font.name = 'OPEN SANS'
    font.size = Pt(36)
    # p.alignment = PP_ALIGN.JUSTIFY
    font.color.rgb = RGBColor(149, 26, 63)
    fill = shape.fill
    fill.background()
    line = shape.line
    line.fill.background()


def add_slide(prs, text_cit, path_img, surname, name, mention_eleve, add_cit=True, add_ima=True):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height
    add_title(shapes, surname, name)
    add_mention(shapes, mention_eleve)
    if add_cit:
        add_citation(shapes, text_cit)
    if add_ima:
        add_image(shapes, path_img, SLIDE_WIDTH, SLIDE_HEIGHT)

    # prs.save('test.pptx')


if __name__ == '__main__':

    csv_path = r'rdd_list.csv'
    #df = pd.read_csv(csv_path, encoding ='utf-8')
    #df.to_csv(csv_path, encoding='utf-8', index=False)
    #with open("rdd_list.csv", encoding="utf16") as csvfile :
      #  print(csvfile)
     #   exemple = csv.reader(csvfile, delimiter=",")
    #csv_path = exemple


    for i in range(1,2): # numéro de la dominante pour
        dominante = str(i) # str(i)
        print(dominante)
        dominante_name = dict_dominante[dominante]
        print(dominante_name)
        students_feat = FramaFormResult(csv_path)
        students_feat_json = students_feat.get_json_dominante(dominante)

        nb_student = students_feat.get_nb_student()
        prs = Presentation()
        prs.slide_height = Inches(9)
        prs.slide_width = Inches(16)
        add_slide_dominante(dominante_name, prs)
        mention = ''
        for i in range(len(students_feat_json)):
            add_cit = True
            add_img = True
            student = students_feat_json[i]

            def get_key(val):
                for key, value in dict_mention.items():
                    if val== value :
                        return key
                return 0

            name = student['nom']
            surname = student['prenom']
            cit = student['citation']
            img_url = student['image']
            mention_eleve = student['dominante_mention']
            #print(mention_eleve1)
            mention_eleve1=dict_mention[mention_eleve]
            print(name)

            if mention_eleve != mention:
                add_slide_dominante(dict_mention[mention_eleve], prs)
                mention = mention_eleve
            img_path = img_url

            #response = requests.get(img_url)
            #img_i= BytesIO(response.content)
            #img_path = Image.open(img_i)


            #f = open('output.jpg','wb')
            #f.write(requests.get(img_url).content)
            #f.close()
            #img_path = f


            # load img into output.png
            if str(img_url) != 'nan':
                add_img =True
                img_path = load_img(img_url, cookie_frama, session)

            else :
                add_img = True
                img_path=r'C:\Users\33652\PycharmProjects\RDD\Picture1.jpg'
            if str(cit) == 'nan':
                add_cit = False
            add_slide(prs, cit, img_path, surname, name, mention_eleve1, add_cit, add_img)
        prs.save('RDD_Slide' + '.pptx')
